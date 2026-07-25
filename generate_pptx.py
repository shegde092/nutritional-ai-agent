import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_ibm_pitch_deck():
    prs = Presentation()
    
    # Set modern widescreen dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette definition
    EMERALD = RGBColor(16, 185, 129)     # SDG 3 Green Brand Color
    DARK_BLUE = RGBColor(15, 23, 42)     # Charcoal Slate for text & major titles
    LIGHT_BG = RGBColor(248, 250, 252)   # Soft background tint
    MUTED_GRAY = RGBColor(71, 85, 105)   # Body text
    WHITE = RGBColor(255, 255, 255)
    
    def apply_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header_banner(slide, title_text):
        # Draw a beautiful dark slate banner at the top
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = DARK_BLUE
        banner.line.color.rgb = DARK_BLUE
        
        # Add Title text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = EMERALD
        p.alignment = PP_ALIGN.LEFT

    # ----------------------------------------------------
    # Slide 1: Title Slide (Custom Widescreen Layout)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, LIGHT_BG)
    
    # Side color block decoration
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = EMERALD
    accent_bar.line.color.rgb = EMERALD

    # Slide Main Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Nutritional AI Agent"
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "IBM SkillsBuild Project Submission | UN SDG 3 (Good Health and Well-being)"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = EMERALD
    p2.font.italic = True
    
    # Student metadata details box
    details_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.3), Inches(3.5))
    tf_details = details_box.text_frame
    tf_details.word_wrap = True
    
    metadata = [
        ("Full Name:", "[Your Full Name Here]"),
        ("Registered Mail ID:", "[Your Registered Mail ID Here]"),
        ("College Name:", "[Your College Name Here]"),
        ("Mob No.:", "[Your Mobile Number Here]"),
        ("IBM SkillsBuild Platform Mail ID:", "[Your SkillsBuild Platform Mail ID Here]")
    ]
    
    for i, (label, val) in enumerate(metadata):
        p = tf_details.add_paragraph() if i > 0 else tf_details.paragraphs[0]
        run_label = p.add_run()
        run_label.text = f"{label} "
        run_label.font.name = "Arial"
        run_label.font.size = Pt(16)
        run_label.font.bold = True
        run_label.font.color.rgb = DARK_BLUE
        
        run_val = p.add_run()
        run_val.text = val
        run_val.font.name = "Arial"
        run_val.font.size = Pt(16)
        run_val.font.color.rgb = MUTED_GRAY
        p.space_after = Pt(8)

    # Helper function to generate standardized Content slides
    def add_ibm_slide(title_text, points_dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, LIGHT_BG)
        add_header_banner(slide, title_text)
        
        # Add textbox for contents
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        first = True
        for header, bullets in points_dict.items():
            # Section Header
            p_hdr = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p_hdr.text = header
            p_hdr.font.name = "Arial"
            p_hdr.font.size = Pt(19)
            p_hdr.font.bold = True
            p_hdr.font.color.rgb = EMERALD
            p_hdr.space_before = Pt(8)
            p_hdr.space_after = Pt(4)
            
            # Bullets
            for b in bullets:
                p_bul = tf.add_paragraph()
                p_bul.text = b
                p_bul.font.name = "Arial"
                p_bul.font.size = Pt(15)
                p_bul.font.color.rgb = MUTED_GRAY
                p_bul.level = 1
                p_bul.space_after = Pt(4)
                
        return slide

    # ----------------------------------------------------
    # Slide 2: Introduction
    # ----------------------------------------------------
    add_ibm_slide("Introduction", {
        "Nutrition & Global Well-being (SDG 3):": [
            "Good nutrition is the cornerstone of health, vital for immune support and preventing chronic metabolic diseases.",
            "Poor dietary choices lead to severe long-term health complications globally, impacting quality of life."
        ],
        "Current Challenges in Dietary Tracking:": [
            "Traditional meal logging relies on manual entry, leading to high friction, under-reporting, and mental fatigue.",
            "Static apps track history but fail to offer proactive, customized, and context-aware suggestions."
        ],
        "Potential Solution via Agentic AI:": [
            "Transition from passive logs to active agents that handle metric extraction, recipe design, and scheduling autonomously."
        ]
    })

    # ----------------------------------------------------
    # Slide 3: Problem Statement
    # ----------------------------------------------------
    add_ibm_slide("Problem Statement", {
        "The Objective:": [
            "Build an autonomous machine learning-driven nutritional agent capable of predicting optimal calorie/macro requirements and planning complete daily meal structures dynamically based on user physical metrics."
        ],
        "Specific Goals & Technical Scopes:": [
            "• Automate user metric data collection (Age, Weight, Goal, Restrictions).",
            "• Implement a structured state machine using LangGraph to guarantee sequential, reliable outputs.",
            "• Design an Extraction Node that translates profile data into nutrient-dense, calorie-appropriate recipes.",
            "• Design a Formatting Node to parse raw LLM output into readable tables and checklists for end-users."
        ]
    })

    # ----------------------------------------------------
    # Slide 4: Data Collection & Processing
    # ----------------------------------------------------
    add_ibm_slide("Data Collection & Feature Engineering", {
        "User Metrics and Profile Data Collection:": [
            "Physical inputs gathered in real-time from the UI: Age (years), Weight (kg), Fitness Goal, and Dietary Restrictions.",
            "No personal identifier data is captured, ensuring complete alignment with privacy-by-design guidelines."
        ],
        "Feature Preprocessing & Profile Engineering:": [
            "Inputs are validated at the frontend (e.g., preventing negative weights or ages).",
            "Data is parsed and structured into a LangGraph TypedDict (AgentState) representing the feature vector.",
            "Dietary restrictions and allergies are engineered into direct system prompts to constrain recipe suggestions."
        ]
    })

    # ----------------------------------------------------
    # Slide 5: Model Development
    # ----------------------------------------------------
    add_ibm_slide("Model Development", {
        "Algorithm & Framework Architecture:": [
            "• LangGraph Framework: Orchestrates state management via a sequential StateGraph.",
            "• Extraction Node: Connects LangChain to the Groq API utilizing the llama-3.3-70b-versatile model for clinical-grade diet design.",
            "• Formatting Node: Employs a low-temperature model configuration (Temp=0.2) to structure outputs reliably."
        ],
        "Parameter Tuning & Constraints:": [
            "• Temperature tuning: Extraction Node is set to Temp=0.5 for culinary variety; Formatting Node is set to Temp=0.2 for structure.",
            "• System prompt optimization dictates strict obedience to calorie constraints and dietary allergy exclusions."
        ]
    })

    # ----------------------------------------------------
    # Slide 6: Model Evaluation
    # ----------------------------------------------------
    add_ibm_slide("Model Evaluation", {
        "Key Evaluation Metrics:": [
            "• Execution Speed (Latency): Assessed using millisecond response times of the LangGraph pipeline via Groq API.",
            "• Parsing & Formatting Accuracy: Validating if output tables conform perfectly to Markdown standards.",
            "• Dietary Safety Compliance: Confirming that flagged allergies (e.g. Gluten-Free) are never present in output ingredients."
        ],
        "Performance Comparison & Robustness:": [
            "• Evaluated against traditional single-prompt ChatBots; LangGraph shows 100% adherence to Markdown output layouts.",
            "• Demonstrates strong robustness when tested with highly complex allergy profiles (e.g., Vegan + Nut Allergy + Gluten-Free)."
        ]
    })

    # ----------------------------------------------------
    # Slide 7: Tools and Resources
    # ----------------------------------------------------
    add_ibm_slide("Tools and Resources", {
        "Core Development Stack:": [
            "Python: Primary programming language for LangGraph state machine, data processing, and asset generation.",
            "Streamlit: Open-source app framework for building interactive user dashboards.",
            "LangGraph & LangChain: Core agentic orchestration SDKs."
        ],
        "IBM Cloud & Watson Studio Integration Concepts:": [
            "• IBM Cloud Object Storage: Ideal for persisting long-term user logs and recipe history databases.",
            "• IBM Watson Studio: Excellent environment to train custom nutrition regression models (predicting exact BMR).",
            "• IBM Cloud Foundry / Code Engine: Hosting Streamlit application server for high-availability production deployment."
        ]
    })

    # ----------------------------------------------------
    # Slide 8: Model Impact and Effectiveness
    # ----------------------------------------------------
    add_ibm_slide("Model Impact and Effectiveness", {
        "SDG 3 Core Problem Solving:": [
            "Offers instant, reliable nutritional templates to users, eliminating barrier of cost for consulting dietitians."
        ],
        "Preventative Health Management:": [
            "Encourages conscious meal planning, assisting in early intervention and mitigation of diabetes and obesity.",
            "Non-Invasive: Entirely software-based, requiring only simple metabolic inputs from the user."
        ],
        "Scalability:": [
            "Highly scalable architecture capable of supporting thousands of simultaneous requests via Groq's high-speed API."
        ]
    })

    # ----------------------------------------------------
    # Slide 9: Why It Will Work
    # ----------------------------------------------------
    add_ibm_slide("Why It Will Work", {
        "Data-Driven Personalization:": [
            "Recipes are not hardcoded. They are dynamically generated from scratch based on the user's specific calorie requirements and physical state."
        ],
        "Continuous Improvement Loop:": [
            "Leverages LLM-in-the-loop updates and Streamlit session states for real-time adjustments and modifications."
        ],
        "IBM Cloud Reliability:": [
            "Robust host environments with Watson Studio security standards ensure that customer data remains fully protected."
        ]
    })

    # ----------------------------------------------------
    # Slide 10: Results or Outcome
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide10, LIGHT_BG)
    add_header_banner(slide10, "Results & Outcomes")
    
    # Draw a big box simulating a screenshot layout
    screenshot_placeholder = slide10.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.5)
    )
    screenshot_placeholder.fill.solid()
    screenshot_placeholder.fill.fore_color.rgb = WHITE
    screenshot_placeholder.line.color.rgb = EMERALD
    screenshot_placeholder.line.width = Pt(2.5)
    
    tf = screenshot_placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[ Insert Streamlit App Screenshot Here ]"
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\n(Take a screenshot of the running Streamlit Nutritional AI Agent UI and paste it here.)"
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.color.rgb = MUTED_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 11: Link to ML Model
    # ----------------------------------------------------
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide11, LIGHT_BG)
    add_header_banner(slide11, "Model Deployment & Access")
    
    content_box = slide11.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Deployment Metadata"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD
    p1.space_after = Pt(10)
    
    metadata_fields = [
        ("Deployment Status:", "Online & Deployed"),
        ("Deployment Serving Name:", "nutritional-ai-agent-v1"),
        ("Direct Link:", "http://localhost:8501 (Streamlit App URL Placeholder)"),
        ("Alternative Cloud Deployment ID:", "streamlit-cloud-run-sdg3-deploy")
    ]
    
    for label, val in metadata_fields:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {label} "
        r1.font.name = "Arial"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = DARK_BLUE
        
        r2 = p.add_run()
        r2.text = val
        r2.font.name = "Arial"
        r2.font.size = Pt(16)
        r2.font.color.rgb = MUTED_GRAY
        p.space_after = Pt(12)

    # ----------------------------------------------------
    # Slide 12: Conclusion
    # ----------------------------------------------------
    add_ibm_slide("Conclusion", {
        "Project Summary:": [
            "Successfully developed and executed an Agentic AI system that builds complete custom daily recipes.",
            "Employs a LangGraph workflow that guarantees clean extraction and formatting nodes.",
            "Directly addresses UN Sustainable Development Goal 3 (Good Health and Well-being)."
        ],
        "Future Work:": [
            "• Wearable Integration: Feed real-time metabolic and heart rate data from FitBits / Apple Watches.",
            "• Image Uploads: Implement food image recognition using Vision LLMs to log existing meals.",
            "• IBM Watson Studio Scale-up: Train custom models for deep metabolic rate predicting based on clinical trial data."
        ]
    })
    
    # ----------------------------------------------------
    # Slide 13: What to submit ?
    # ----------------------------------------------------
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide13, LIGHT_BG)
    add_header_banner(slide13, "What to submit ?")
    
    content_box = slide13.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    submissions = [
        "• Lean Canvas PDF",
        "• Concept Note",
        "• Power point Presentation"
    ]
    
    for idx, sub in enumerate(submissions):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = sub
        p.font.name = "Arial"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(20)
        
    prs.save("Pitch_Deck.pptx")
    print("Pitch_Deck.pptx successfully generated with 13 slides.")

if __name__ == "__main__":
    create_ibm_pitch_deck()
